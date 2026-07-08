#!/usr/bin/env python
"""Build AI-900 WSQ courseware from the lab Markdown source.

Outputs:
- Microsoft Azure AI Fundamentals AI-900-v1.pptx
- Microsoft Azure AI Fundamentals AI-900-v1.pdf
- LG-Microsoft Azure AI Fundamentals AI-900.docx
- LG-Microsoft Azure AI Fundamentals AI-900.pdf
- LG-Microsoft Azure AI Fundamentals AI-900.md
- LP-Microsoft Azure AI Fundamentals AI-900.docx
- LP-Microsoft Azure AI Fundamentals AI-900.pdf
"""

from __future__ import annotations

import os
import re
import html
import textwrap
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DInches
from docx.shared import Pt, RGBColor as DRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4, landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    ListFlowable,
    ListItem,
    PageBreak,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
COURSEWARE = ROOT / "courseware"
LABS_DIR = ROOT / ".tmp-ai900-repo" / "labs"
REFERENCE_PPT = ROOT / "Resources" / "WSQ - Master Trainer Slides - Microsoft Azure Data Fundamentals (DP-900) - v10.pptx"
EXAMPLE_PPT = ROOT / ".tmp-example-repo" / "courseware" / "Agentic AI Automation with n8n-v45.pptx"
ASSET_DIR = COURSEWARE / "assets" / "imported-reference-media"
LG_DIAGRAM_DIR = COURSEWARE / "assets" / "lg-diagrams"
WSQ_LOGO = ASSET_DIR / "dp900-image1.png"
TERTIARY_MARK = ASSET_DIR / "dp900-image2.png"
TERTIARY_WORDMARK = ASSET_DIR / "dp900-image3.png"
COURSE_BADGE = COURSEWARE / "assets" / "ai900-course-badge.png"

TITLE = "Microsoft Azure AI Fundamentals AI-900"
TITLE_WITH_PARENS = "Microsoft Azure AI Fundamentals (AI-900)"
CODE = "TGS-2023021100"
TSC_TITLE = "Artificial Intelligence Application"
TSC_CODE = "ICT-DIT-4016-1.1"
VERSION = "v1"
DOC_VERSION = "1.0"
ORG = "Tertiary Infotech Academy Pte Ltd"
UEN = "UEN 201200696W"
URL = "www.tertiarycourses.com.sg"
LAB_REPO = "github.com/tertiarycourses/TGS-2023021100---Microsoft-Azure-AI-Fundamentals-AI-900-"
DATE = "8 July 2026"

BLUE = RGBColor(0x1F, 0x6F, 0xEB)
TEAL = RGBColor(0x10, 0xB9, 0x81)
NAVY = RGBColor(0x0B, 0x12, 0x20)
INK = RGBColor(0x16, 0x1B, 0x26)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x12, 0x7A, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"


@dataclass
class Lab:
    num: int
    title: str
    filename: str
    objectives: list[str] = field(default_factory=list)
    scenario: str = ""
    steps_raw: list[str] = field(default_factory=list)
    step_titles: list[str] = field(default_factory=list)
    validation: str = ""
    checkpoint: list[str] = field(default_factory=list)
    focus: str = ""


def ensure_dirs() -> None:
    COURSEWARE.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    if not COURSE_BADGE.exists():
        from PIL import Image, ImageDraw, ImageFont

        img = Image.new("RGBA", (560, 360), (255, 255, 255, 0))
        d = ImageDraw.Draw(img)
        blue = (31, 111, 235, 255)
        navy = (17, 24, 39, 255)
        teal = (16, 185, 129, 255)
        d.ellipse((130, 20, 430, 320), outline=(210, 218, 230, 255), width=8, fill=(255, 255, 255, 255))
        d.rounded_rectangle((95, 134, 465, 210), radius=10, fill=blue)
        try:
            font_big = ImageFont.truetype("arialbd.ttf", 64)
            font_med = ImageFont.truetype("arialbd.ttf", 30)
            font_small = ImageFont.truetype("arial.ttf", 21)
        except Exception:
            font_big = font_med = font_small = None
        def center(text, y, font, fill):
            box = d.textbbox((0, 0), text, font=font)
            d.text(((560 - (box[2] - box[0])) / 2, y), text, font=font, fill=fill)
        center("Microsoft Azure", 72, font_small, navy)
        center("AI-900", 137, font_big, (255, 255, 255, 255))
        center("AI Fundamentals", 226, font_med, navy)
        d.arc((170, 48, 390, 300), 205, 335, fill=teal, width=6)
        img.save(COURSE_BADGE)


def extract_reference_media() -> None:
    """Import available visual media from the reference/example PPTs."""
    for ppt in [REFERENCE_PPT, EXAMPLE_PPT]:
        if not ppt.exists():
            continue
        prefix = "dp900" if "DP-900" in ppt.name else "example"
        with zipfile.ZipFile(ppt) as zf:
            for name in zf.namelist():
                if not name.startswith("ppt/media/"):
                    continue
                ext = Path(name).suffix.lower()
                if ext not in [".png", ".jpg", ".jpeg"]:
                    continue
                data = zf.read(name)
                if len(data) < 2048:
                    continue
                out = ASSET_DIR / f"{prefix}-{Path(name).name}"
                if not out.exists():
                    out.write_bytes(data)


def strip_md(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


def collect_section(lines: list[str], heading: str) -> list[str]:
    start = None
    for i, line in enumerate(lines):
        if line.strip().lower() == f"## {heading}".lower():
            start = i + 1
            break
    if start is None:
        return []
    end = len(lines)
    for j in range(start, len(lines)):
        if lines[j].startswith("## "):
            end = j
            break
    return lines[start:end]


def clean_block(lines: list[str]) -> str:
    return "\n".join(lines).strip()


def parse_bullets(lines: list[str]) -> list[str]:
    out = []
    for line in lines:
        s = line.strip()
        if s.startswith("- "):
            out.append(strip_md(s[2:]))
        elif re.match(r"^\d+\.\s+", s):
            out.append(strip_md(re.sub(r"^\d+\.\s+", "", s)))
    return out


def parse_labs() -> list[Lab]:
    labs = []
    for path in sorted(LABS_DIR.glob("lab-*.md")):
        text = path.read_text(encoding="utf-8")
        lines = text.splitlines()
        title_line = next((x for x in lines if x.startswith("# ")), "# Lab")
        m = re.match(r"# Lab\s+(\d+)\s*-\s*(.+)", title_line)
        if not m:
            continue
        lab = Lab(num=int(m.group(1)), title=strip_md(m.group(2)), filename=path.name)
        lab.objectives = parse_bullets(collect_section(lines, "Objectives"))
        lab.scenario = strip_md(clean_block(collect_section(lines, "Scenario")))
        lab.steps_raw = collect_section(lines, "Steps")
        lab.step_titles = [
            strip_md(re.sub(r"^###\s+", "", x))
            for x in lab.steps_raw
            if x.startswith("### ")
        ]
        lab.validation = strip_md(clean_block(collect_section(lines, "Validation")))
        lab.checkpoint = parse_bullets(collect_section(lines, "Checkpoint Questions"))
        lab.focus = strip_md(clean_block(collect_section(lines, "Exam Focus") or collect_section(lines, "Course Focus")))
        labs.append(lab)
    return labs


def ensure_lg_diagrams(labs: list[Lab]) -> None:
    from PIL import Image, ImageDraw, ImageFont

    LG_DIAGRAM_DIR.mkdir(parents=True, exist_ok=True)
    palette = {
        "blue": (31, 111, 235, 255),
        "teal": (16, 185, 129, 255),
        "amber": (245, 158, 11, 255),
        "purple": (124, 58, 237, 255),
        "ink": (17, 24, 39, 255),
        "grey": (85, 91, 102, 255),
        "line": (203, 213, 225, 255),
        "light": (248, 250, 252, 255),
        "white": (255, 255, 255, 255),
    }

    def font(size: int, bold: bool = False):
        names = ["arialbd.ttf", "arial.ttf"] if bold else ["arial.ttf", "calibri.ttf"]
        for name in names:
            try:
                return ImageFont.truetype(name, size)
            except Exception:
                continue
        return ImageFont.load_default()

    f_title = font(44, True)
    f_head = font(26, True)
    f_body = font(22)
    f_small = font(18)
    f_num = font(24, True)

    def wrap(text: str, fnt, max_width: int) -> list[str]:
        words = text.split()
        lines: list[str] = []
        current = ""
        scratch = Image.new("RGB", (10, 10))
        d = ImageDraw.Draw(scratch)
        for word in words:
            candidate = f"{current} {word}".strip()
            if d.textbbox((0, 0), candidate, font=fnt)[2] <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = word
        if current:
            lines.append(current)
        return lines

    def draw_wrapped(d: ImageDraw.ImageDraw, text: str, xy: tuple[int, int], fnt, fill, max_width: int, line_gap: int = 5) -> int:
        x, y = xy
        for line in wrap(text, fnt, max_width):
            d.text((x, y), line, font=fnt, fill=fill)
            box = d.textbbox((x, y), line, font=fnt)
            y = box[3] + line_gap
        return y

    def save_learning_flow() -> None:
        img = Image.new("RGBA", (1500, 760), palette["white"])
        d = ImageDraw.Draw(img)
        d.rectangle((0, 0, 1500, 12), fill=palette["blue"])
        d.text((70, 54), "Learner Guide Flow", font=f_title, fill=palette["ink"])
        d.text((70, 112), "Each topic moves from concept briefing to hands-on evidence and assessment readiness.", font=f_body, fill=palette["grey"])
        steps = [
            ("1", "Concept", "Recognize the AI workload and responsible AI concern.", palette["blue"]),
            ("2", "Guided Lab", "Follow the Azure AI lab steps and capture evidence.", palette["teal"]),
            ("3", "Validation", "Check output quality, service fit, and resource status.", palette["purple"]),
            ("4", "Checkpoint", "Answer review questions using lab evidence.", palette["amber"]),
            ("5", "Review", "Map scenario, service, control, and cleanup action.", palette["blue"]),
        ]
        x, y, w, h, gap = 70, 245, 252, 285, 27
        for idx, (num, title, body, col) in enumerate(steps):
            left = x + idx * (w + gap)
            d.rounded_rectangle((left, y, left + w, y + h), radius=22, fill=palette["light"], outline=col, width=3)
            d.ellipse((left + 22, y + 24, left + 82, y + 84), fill=col)
            d.text((left + 44, y + 39), num, font=f_num, fill=palette["white"], anchor="mm")
            d.text((left + 24, y + 116), title, font=f_head, fill=col)
            draw_wrapped(d, body, (left + 24, y + 162), f_body, palette["ink"], w - 48)
            if idx < len(steps) - 1:
                d.text((left + w + 11, y + 130), ">", font=font(38, True), fill=palette["grey"])
        d.text((70, 650), "Use this flow for every lab in the guide: read the scenario, complete the steps, validate evidence, and answer the checkpoint.", font=f_body, fill=palette["grey"])
        img.save(LG_DIAGRAM_DIR / "course-learning-flow.png")

    def save_service_map() -> None:
        img = Image.new("RGBA", (1500, 860), palette["white"])
        d = ImageDraw.Draw(img)
        d.rectangle((0, 0, 1500, 12), fill=palette["teal"])
        d.text((70, 54), "Azure AI Service Map", font=f_title, fill=palette["ink"])
        d.text((70, 112), "Start from the input and output, then choose the closest Azure AI capability.", font=f_body, fill=palette["grey"])
        cards = [
            (70, 205, "Machine Learning", "Azure Machine Learning", "Predict values, classify, cluster", palette["blue"]),
            (70, 375, "Vision", "Azure AI Vision", "Analyze images, OCR, detect objects", palette["teal"]),
            (70, 545, "Language", "Azure AI Language", "Sentiment, entities, key phrases", palette["amber"]),
            (980, 205, "Speech & Translation", "Azure AI Speech / Translator", "Transcribe, synthesize, translate", palette["blue"]),
            (980, 375, "Documents & Search", "Document Intelligence / AI Search", "Extract forms, index knowledge", palette["teal"]),
            (980, 545, "Generative AI", "Azure AI Foundry / Azure OpenAI", "Prompt, ground, review outputs", palette["amber"]),
        ]
        d.rounded_rectangle((560, 300, 940, 430), radius=28, fill=palette["blue"])
        d.text((750, 348), "Azure AI", font=f_title, fill=palette["white"], anchor="mm")
        d.text((750, 392), "service selection hub", font=f_small, fill=palette["white"], anchor="mm")
        d.rounded_rectangle((580, 475, 920, 565), radius=16, fill=palette["light"], outline=palette["line"], width=2)
        d.text((750, 520), "Match scenario to service", font=f_body, fill=palette["ink"], anchor="mm")
        for x, y, head, service, cue, col in cards:
            d.rounded_rectangle((x, y, x + 390, y + 120), radius=14, fill=palette["white"], outline=col, width=3)
            d.rectangle((x, y, x + 12, y + 120), fill=col)
            d.text((x + 30, y + 22), head, font=f_head, fill=col)
            d.text((x + 30, y + 58), service, font=f_body, fill=palette["ink"])
            d.text((x + 30, y + 88), cue, font=f_small, fill=palette["grey"])
        d.text((130, 775), "Lab use: identify workload -> choose service -> validate output -> name responsible AI control.", font=f_body, fill=palette["grey"])
        img.save(LG_DIAGRAM_DIR / "azure-ai-service-map.png")

    def save_lab_flow(lab: Lab) -> None:
        img = Image.new("RGBA", (1500, 760), palette["white"])
        d = ImageDraw.Draw(img)
        d.rectangle((0, 0, 1500, 12), fill=palette["blue"])
        d.text((70, 48), f"Lab {lab.num:02d} Workflow", font=f_title, fill=palette["ink"])
        draw_wrapped(d, lab.title, (70, 108), f_body, palette["grey"], 1120)
        steps = lab.step_titles[:5] or ["Read scenario", "Complete guided steps", "Validate output", "Answer checkpoint"]
        if len(lab.step_titles) > 5:
            steps[-1] = "Validate and review"
        count = len(steps)
        w = 250 if count >= 5 else 300
        gap = 24
        total = count * w + (count - 1) * gap
        x0 = (1500 - total) // 2
        y = 260
        colors_cycle = [palette["blue"], palette["teal"], palette["purple"], palette["amber"], palette["blue"]]
        for idx, step in enumerate(steps):
            x = x0 + idx * (w + gap)
            col = colors_cycle[idx % len(colors_cycle)]
            d.rounded_rectangle((x, y, x + w, y + 205), radius=18, fill=palette["light"], outline=col, width=3)
            d.ellipse((x + 20, y + 22, x + 72, y + 74), fill=col)
            d.text((x + 46, y + 48), str(idx + 1), font=f_num, fill=palette["white"], anchor="mm")
            draw_wrapped(d, step, (x + 24, y + 98), f_body, palette["ink"], w - 48)
            if idx < count - 1:
                d.text((x + w + 8, y + 78), ">", font=font(34, True), fill=palette["grey"])
        d.rounded_rectangle((100, 560, 1400, 650), radius=18, fill=palette["light"], outline=palette["line"], width=2)
        draw_wrapped(d, f"Validation: {lab.validation}", (130, 585), f_small, palette["ink"], 1240)
        img.save(LG_DIAGRAM_DIR / f"lab-{lab.num:02d}-workflow.png")

    save_learning_flow()
    save_service_map()
    for lab in labs:
        save_lab_flow(lab)


class SlideBuilder:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.slide_index: list[dict[str, str | int]] = []
        self.page = 0

    def _slide(self, title: str, section: str) -> object:
        slide = self.prs.slides.add_slide(self.blank)
        self.slide_index.append({"slide": len(self.slide_index) + 1, "title": title, "section": section})
        return slide

    def rect(self, slide, x, y, w, h, fill, line=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = PPt(1)
        sh.shadow.inherit = False
        return sh

    def rrect(self, slide, x, y, w, h, fill, line=LINE):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            sh.adjustments[0] = 0.05
        except Exception:
            pass
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = PPt(1)
        sh.shadow.inherit = False
        return sh

    def oval(self, slide, x, y, w, h, fill):
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def text(self, slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = PPt(0)
        tf.margin_top = tf.margin_bottom = PPt(0)
        if runs and not isinstance(runs[0], list):
            runs = [runs]
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = PPt(space)
            p.space_before = PPt(0)
            for item in para:
                if len(item) == 4:
                    t, sz, col, bold = item
                    font = FONT
                else:
                    t, sz, col, bold, font = item
                r = p.add_run()
                r.text = str(t)
                r.font.size = PPt(sz)
                r.font.bold = bool(bold)
                r.font.color.rgb = col
                r.font.name = font
        return tb

    def bullets(self, slide, x, y, w, h, items, size=18, color=INK, gap=8, marker=BLUE):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = PPt(0)
        tf.margin_top = tf.margin_bottom = PPt(0)
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PPt(gap)
            p.space_before = PPt(0)
            rm = p.add_run()
            rm.text = "-  "
            rm.font.size = PPt(size)
            rm.font.bold = True
            rm.font.color.rgb = marker
            rm.font.name = FONT
            r = p.add_run()
            r.text = str(item)
            r.font.size = PPt(size)
            r.font.color.rgb = color
            r.font.name = FONT
        return tb

    def footer(self, slide, dark=False):
        c = WHITE if dark else GREY
        slide_no = len(self.slide_index)
        self.text(slide, Inches(0.55), Inches(7.08), Inches(5.8), Inches(0.3), [[(f"{TITLE_WITH_PARENS} - {CODE}", 8.5, c, False)]])
        self.text(slide, Inches(5.1), Inches(7.08), Inches(4.3), Inches(0.3), [[("(c) 2026 Tertiary Infotech Academy Pte Ltd", 8.5, c, False)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(11.9), Inches(7.08), Inches(0.9), Inches(0.3), [[(str(slide_no), 9, c, True)]], align=PP_ALIGN.RIGHT)

    def head(self, slide, title, kicker=None, accent=BLUE):
        self.rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, WHITE)
        self.rect(slide, Inches(0.55), Inches(0.60), Inches(0.14), Inches(0.66), accent)
        if kicker:
            self.text(slide, Inches(0.86), Inches(0.50), Inches(8), Inches(0.25), [[(kicker.upper(), 9.5, accent, True)]])
            y = Inches(0.78)
        else:
            y = Inches(0.60)
        self.text(slide, Inches(0.85), y, Inches(10.8), Inches(0.70), [[(title, 31, INK, True)]])

    def ref_head(self, slide, title, kicker, accent=BLUE, grey=False):
        self.rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, WHITE)
        bar = RGBColor(0x62, 0x6B, 0x7A) if grey else accent
        self.rect(slide, Inches(0.65), Inches(0.62), Inches(0.10), Inches(0.70), bar)
        self.text(slide, Inches(0.92), Inches(0.55), Inches(8.5), Inches(0.32), [[(kicker.upper(), 14, RGBColor(0x5B, 0x63, 0x72) if grey else accent, True)]])
        self.text(slide, Inches(0.92), Inches(0.92), Inches(11.0), Inches(0.72), [[(title, 29, INK, True)]])
        self.rect(slide, Inches(0.9), Inches(1.64), Inches(11.9), Inches(0.03), RGBColor(0xD9, 0xE1, 0xEA))

    def cover(self):
        slide = self._slide(TITLE_WITH_PARENS, "Admin")
        self.rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, WHITE)
        self.rect(slide, 0, Inches(0.10), self.prs.slide_width, Inches(0.20), BLUE)
        self.rect(slide, 0, Inches(7.22), self.prs.slide_width, Inches(0.22), TEAL)
        if TERTIARY_MARK.exists():
            slide.shapes.add_picture(str(TERTIARY_MARK), Inches(0.88), Inches(0.78), height=Inches(1.00))
        if COURSE_BADGE.exists():
            self.rrect(slide, Inches(10.10), Inches(0.82), Inches(1.52), Inches(0.98), RGBColor(0xF5, 0x9E, 0x0B), line=None)
            self.text(slide, Inches(10.10), Inches(1.03), Inches(1.52), Inches(0.26), [[("AI-900", 18, WHITE, True)]], align=PP_ALIGN.CENTER)
            self.text(slide, Inches(10.10), Inches(1.42), Inches(1.52), Inches(0.18), [[("AZURE AI", 8, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(0.98), Inches(2.40), Inches(4.2), Inches(0.30), [[("COURSE SLIDES  ·  WSQ", 15, BLUE, True)]])
        self.text(slide, Inches(0.98), Inches(2.95), Inches(10.90), Inches(1.08), [[(TITLE_WITH_PARENS, 34, INK, True)]])
        self.rect(slide, Inches(0.98), Inches(4.56), Inches(2.22), Inches(0.055), TEAL)
        self.text(slide, Inches(0.98), Inches(4.94), Inches(7.0), Inches(0.28), [[(f"WSQ Course Code: {CODE}", 13.5, GREY, False)]])
        self.text(slide, Inches(0.98), Inches(5.25), Inches(7.0), Inches(0.25), [[("Trainer: Dr. Alfred Ang", 11.5, GREY, False)]])
        self.text(slide, Inches(0.98), Inches(5.55), Inches(9.0), Inches(0.25), [[(f"Conducted by {ORG}  ·  UEN 201200696W", 11.5, GREY, False)]])
        self.text(slide, Inches(0.98), Inches(6.28), Inches(4.2), Inches(0.22), [[(f"Version {DOC_VERSION}  ·  {DATE}", 9.5, GREY, False)]])
        self.text(slide, Inches(0.98), Inches(6.62), Inches(7.5), Inches(0.18), [[("(c) 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg", 7.2, GREY, False)]])

    def trainer_template_slide(self):
        slide = self._slide("About the Trainer", "Admin")
        self.ref_head(slide, "About the Trainer", "Your Trainer - General", grey=True)
        self.rrect(slide, Inches(0.98), Inches(2.05), Inches(3.65), Inches(4.65), LIGHT, line=None)
        self.rect(slide, Inches(0.98), Inches(2.05), Inches(3.65), Inches(0.11), RGBColor(0x62, 0x6B, 0x7A))
        self.oval(slide, Inches(1.88), Inches(2.78), Inches(1.72), Inches(1.72), RGBColor(0x62, 0x6B, 0x7A))
        self.text(slide, Inches(1.88), Inches(3.25), Inches(1.72), Inches(0.38), [[("?", 30, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(1.45), Inches(4.68), Inches(2.7), Inches(0.4), [[("Your Trainer", 20, INK, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(1.25), Inches(5.25), Inches(3.1), Inches(0.65), [[("General Trainer template -\nto be completed by the trainer", 12, GREY, True)]], align=PP_ALIGN.CENTER)
        fields = [("NAME", BLUE), ("TITLE / DESIGNATION", TEAL), ("QUALIFICATIONS", RGBColor(0x7C, 0x3A, 0xED)), ("AREAS OF EXPERTISE", RGBColor(0xF5, 0x9E, 0x0B)), ("TRAINING & INDUSTRY EXPERIENCE", BLUE), ("CONTACT", TEAL)]
        y = Inches(2.05)
        for label, col in fields:
            self.rrect(slide, Inches(5.02), y, Inches(7.65), Inches(0.63), LIGHT, line=None)
            self.rect(slide, Inches(5.02), y, Inches(0.09), Inches(0.63), col)
            self.text(slide, Inches(5.33), y + Inches(0.16), Inches(4.8), Inches(0.22), [[(label, 10.5, col, True)]])
            self.rect(slide, Inches(5.32), y + Inches(0.46), Inches(4.45), Inches(0.015), RGBColor(0xC9, 0xD6, 0xE6))
            y += Inches(0.80)
        self.footer(slide)

    def trainer_profile_slide(self):
        slide = self._slide("About the Trainer", "Admin")
        self.ref_head(slide, "About the Trainer", "Your Trainer", BLUE)
        self.rrect(slide, Inches(1.02), Inches(2.08), Inches(3.65), Inches(4.62), LIGHT, line=LINE)
        self.rect(slide, Inches(1.02), Inches(2.08), Inches(3.65), Inches(0.11), BLUE)
        self.oval(slide, Inches(1.9), Inches(2.62), Inches(1.72), Inches(1.72), BLUE)
        self.text(slide, Inches(1.9), Inches(3.09), Inches(1.72), Inches(0.38), [[("AA", 30, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(1.35), Inches(4.75), Inches(3.0), Inches(0.35), [[("Dr. Alfred Ang", 19, INK, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(1.35), Inches(5.35), Inches(3.0), Inches(0.55), [[("Principal Trainer\nTertiary Infotech Academy Pte. Ltd.", 11.5, GREY, True)]], align=PP_ALIGN.CENTER)
        rows = [
            ("ROLE", "Principal Trainer, Tertiary Infotech Academy Pte. Ltd.", BLUE),
            ("QUALIFICATIONS", "PhD - specialises in AI, automation and software engineering.", TEAL),
            ("DELIVERS", "WSQ courses on AI agents, automation, cloud, data and app development.", RGBColor(0x7C, 0x3A, 0xED)),
            ("FOUNDER", "Founder and lead instructor at Tertiary Infotech / Tertiary Courses.", RGBColor(0xF5, 0x9E, 0x0B)),
        ]
        y = Inches(2.08)
        for label, body, col in rows:
            self.rrect(slide, Inches(5.05), y, Inches(7.65), Inches(0.98), LIGHT, line=LINE)
            self.rect(slide, Inches(5.05), y, Inches(0.09), Inches(0.98), col)
            self.text(slide, Inches(5.36), y + Inches(0.22), Inches(2.0), Inches(0.22), [[(label, 10.5, col, True)]])
            self.text(slide, Inches(5.36), y + Inches(0.50), Inches(6.8), Inches(0.32), [[(body, 14, INK, False)]])
            y += Inches(1.22)
        self.footer(slide)

    def ground_rules_slide(self):
        slide = self._slide("Ground Rules", "Admin")
        self.ref_head(slide, "Ground Rules", "Housekeeping", BLUE)
        rules = [
            ("1", "Set your mobile phone to silent mode.", BLUE),
            ("2", "Participate actively - no question is too small.", TEAL),
            ("3", "Mutual respect: agree to disagree.", RGBColor(0x7C, 0x3A, 0xED)),
            ("4", "One conversation at a time.", RGBColor(0xF5, 0x9E, 0x0B)),
            ("5", "Be punctual; return from breaks on time.", RGBColor(0xE1, 0x1D, 0x48)),
            ("6", "75% attendance is required.", RGBColor(0x0E, 0xA5, 0xE9)),
        ]
        for idx, (num, body, col) in enumerate(rules):
            x = Inches(0.68 if idx % 2 == 0 else 6.88)
            y = Inches(1.95 + (idx // 2) * 1.8)
            self.rrect(slide, x, y, Inches(5.9), Inches(1.25), LIGHT, line=LINE)
            self.rect(slide, x, y, Inches(0.09), Inches(1.25), col)
            self.oval(slide, x + Inches(0.32), y + Inches(0.36), Inches(0.62), Inches(0.62), col)
            self.text(slide, x + Inches(0.32), y + Inches(0.49), Inches(0.62), Inches(0.24), [[(num, 14, WHITE, True)]], align=PP_ALIGN.CENTER)
            self.text(slide, x + Inches(1.15), y + Inches(0.40), Inches(4.45), Inches(0.45), [[(body, 17, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
        self.footer(slide)

    def admin_slide(self, title, kicker, items):
        slide = self._slide(title, "Admin")
        self.head(slide, title, kicker)
        self.bullets(slide, Inches(0.95), Inches(1.75), Inches(11.4), Inches(4.9), items, size=19)
        self.footer(slide)

    def two_col(self, title, kicker, left_title, left_items, right_title, right_items):
        slide = self._slide(title, "Admin")
        self.head(slide, title, kicker)
        self.rrect(slide, Inches(0.85), Inches(1.85), Inches(5.65), Inches(4.85), LIGHT)
        self.rrect(slide, Inches(6.85), Inches(1.85), Inches(5.65), Inches(4.85), LIGHT)
        self.text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.35), [[(left_title, 16, BLUE, True)]])
        self.text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.35), [[(right_title, 16, TEAL, True)]])
        self.bullets(slide, Inches(1.1), Inches(2.65), Inches(5.0), Inches(3.7), left_items, size=15)
        self.bullets(slide, Inches(7.1), Inches(2.65), Inches(5.0), Inches(3.7), right_items, size=15, marker=TEAL)
        self.footer(slide)

    def section(self, kicker, title, sub=""):
        slide = self._slide(title, "Section")
        self.rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, WHITE)
        self.rect(slide, 0, 0, Inches(0.28), self.prs.slide_height, BLUE)
        self.rect(slide, Inches(0.85), Inches(2.45), Inches(0.14), Inches(1.85), TEAL)
        self.text(slide, Inches(1.25), Inches(2.5), Inches(11), Inches(0.5), [[(kicker.upper(), 18, BLUE, True)]])
        self.text(slide, Inches(1.25), Inches(3.0), Inches(11.4), Inches(1.3), [[(title, 38, INK, True)]])
        if sub:
            self.text(slide, Inches(1.27), Inches(4.35), Inches(11), Inches(0.8), [[(sub, 16, GREY, False)]])
        self.footer(slide)

    def content_slide(self, title, kicker, items, size=18, accent=BLUE):
        slide = self._slide(title, kicker)
        self.head(slide, title, kicker, accent)
        self.bullets(slide, Inches(0.9), Inches(1.75), Inches(11.5), Inches(5.0), items, size=size, marker=accent)
        self.footer(slide)

    def diagram_slide(self):
        slide = self._slide("Azure AI services match workloads to capabilities", "Concept")
        self.head(slide, "Azure AI services match workloads to capabilities", "CONCEPT", TEAL)
        self.text(slide, Inches(0.9), Inches(1.48), Inches(11.4), Inches(0.34), [[("Start from the workload, then select the Azure AI capability that best fits the input, output, and control needs.", 13.5, GREY, False)]])

        self.rrect(slide, Inches(4.85), Inches(2.38), Inches(3.15), Inches(1.35), BLUE, line=None)
        self.text(slide, Inches(5.15), Inches(2.68), Inches(2.55), Inches(0.28), [[("Azure AI", 22, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(5.15), Inches(3.12), Inches(2.55), Inches(0.22), [[("Service selection hub", 10.5, WHITE, False)]], align=PP_ALIGN.CENTER)
        self.rrect(slide, Inches(4.85), Inches(4.02), Inches(3.15), Inches(0.88), LIGHT, line=LINE)
        self.text(slide, Inches(5.15), Inches(4.25), Inches(2.55), Inches(0.24), [[("Match scenario to service", 12, INK, True)]], align=PP_ALIGN.CENTER)

        cards = [
            (0.9, 2.05, "Machine Learning", "Azure Machine Learning", "Predict values, classify, cluster", BLUE),
            (0.9, 3.58, "Vision", "Azure AI Vision", "Analyze images, OCR, detect objects", TEAL),
            (0.9, 5.11, "Language", "Azure AI Language", "Sentiment, entities, key phrases", AMBER),
            (8.55, 2.05, "Speech & Translation", "Azure AI Speech / Translator", "Transcribe, synthesize, translate", BLUE),
            (8.55, 3.58, "Documents & Search", "Document Intelligence / AI Search", "Extract forms, index knowledge", TEAL),
            (8.55, 5.11, "Generative AI", "Azure AI Foundry / Azure OpenAI", "Prompt, ground, review outputs", AMBER),
        ]
        for x, y, heading, service, cue, col in cards:
            self.rrect(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.10), WHITE, line=col)
            self.rect(slide, Inches(x), Inches(y), Inches(0.08), Inches(1.10), col)
            self.text(slide, Inches(x + 0.22), Inches(y + 0.17), Inches(3.05), Inches(0.20), [[(heading, 12.5, col, True)]])
            self.text(slide, Inches(x + 0.22), Inches(y + 0.45), Inches(3.05), Inches(0.20), [[(service, 10.5, INK, True)]])
            self.text(slide, Inches(x + 0.22), Inches(y + 0.72), Inches(3.05), Inches(0.18), [[(cue, 9.3, GREY, False)]])

        self.text(slide, Inches(1.0), Inches(6.55), Inches(11.0), Inches(0.28), [[("Lab connection: learners repeatedly identify the workload, choose the service, validate the result, and name a responsible AI control.", 12.5, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(slide)

    def skills_framework_diagram(self):
        slide = self._slide("Skills Framework links abilities to labs and assessment", "Admin")
        self.head(slide, "Skills Framework links abilities to labs and assessment", "SKILLS FRAMEWORK", BLUE)
        cols = [
            ("Knowledge", ["AI workloads", "Responsible AI", "Azure AI services", "Governance controls"], BLUE),
            ("Abilities", ["Map scenarios", "Select services", "Validate outputs", "Explain risks"], TEAL),
            ("Practice", ["Labs 01-10", "Checkpoints", "Capstone service map", "Cleanup review"], AMBER),
            ("Assessment", ["Open-book MCQ", "Competency evidence", "Attendance", "TRAQOM"], BLUE),
        ]
        x0 = Inches(0.85)
        y = Inches(1.95)
        w = Inches(2.75)
        gap = Inches(0.28)
        for idx, (head, items, col) in enumerate(cols):
            x = x0 + idx * (w + gap)
            self.rrect(slide, x, y, w, Inches(3.9), LIGHT, line=col)
            self.text(slide, x + Inches(0.18), y + Inches(0.22), w - Inches(0.36), Inches(0.34), [[(head, 15, col, True)]], align=PP_ALIGN.CENTER)
            self.bullets(slide, x + Inches(0.22), y + Inches(0.82), w - Inches(0.44), Inches(2.6), items, size=11.5, gap=6, marker=col)
            if idx < len(cols) - 1:
                self.text(slide, x + w + Inches(0.02), y + Inches(1.7), gap - Inches(0.04), Inches(0.45), [[(">", 24, BLUE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(0.9), Inches(6.18), Inches(11.4), Inches(0.38), [[("Alignment rule: every course outcome is practiced in the labs and referenced in the lesson plan slide numbers.", 13, GREY, False)]])
        self.footer(slide)

    def tsc_detail_slide(self, title, kicker, rows, caption):
        slide = self._slide(title, "Admin")
        self.head(slide, title, kicker, BLUE)
        self.text(slide, Inches(0.9), Inches(1.52), Inches(11.4), Inches(0.3), [[(caption, 13, GREY, False)]])
        y = Inches(2.0)
        for idx, (code, text) in enumerate(rows, 1):
            fill = WHITE if idx % 2 else LIGHT
            self.rrect(slide, Inches(0.9), y, Inches(11.45), Inches(0.52), fill, line=LINE)
            self.text(slide, Inches(1.1), y + Inches(0.12), Inches(0.75), Inches(0.22), [[(code, 11, BLUE, True)]], align=PP_ALIGN.CENTER)
            self.text(slide, Inches(1.95), y + Inches(0.10), Inches(10.0), Inches(0.28), [[(text, 11.5, INK, False)]])
            y += Inches(0.6)
            if y > Inches(6.2):
                break
        self.footer(slide)

    def assessment_flow(self):
        slide = self._slide("Assessment Flow", "Admin")
        self.ref_head(slide, "Assessment Flow Reminder", "Assessment", BLUE)
        steps = [
            ("1", "TRAQOM digital attendance"),
            ("2", "Assessment digital attendance"),
            ("3", "Sit written / MCQ assessment"),
            ("4", "Submit answers on LMS/TMS"),
            ("5", "Sign Assessment Summary Record"),
        ]
        gap = Inches(0.25)
        box_w = (Inches(11.5) - gap * 4) // 5
        y = Inches(2.45)
        for idx, (num, label) in enumerate(steps):
            x = Inches(0.9) + idx * (box_w + gap)
            col = [BLUE, TEAL, RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xF5, 0x9E, 0x0B), BLUE][idx]
            self.rrect(slide, x, y, box_w, Inches(1.9), WHITE, line=LINE)
            self.oval(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.50), Inches(0.50), col)
            self.text(slide, x + Inches(0.18), y + Inches(0.32), Inches(0.50), Inches(0.22), [[(num, 12, WHITE, True)]], align=PP_ALIGN.CENTER)
            self.text(slide, x + Inches(0.26), y + Inches(0.83), box_w - Inches(0.52), Inches(0.56), [[(label, 11.5, INK, True)]], align=PP_ALIGN.CENTER)
            if idx < len(steps) - 1:
                self.text(slide, x + box_w + Inches(0.02), y + Inches(0.83), gap - Inches(0.04), Inches(0.35), [[(">", 20, GREY, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(2.35), Inches(5.35), Inches(8.7), Inches(0.45), [[("The TRAQOM QR code and LMS/TMS submission are handled through the LMS/TMS portal.", 15.5, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(slide)

    def practice_exam_slide(self):
        slide = self._slide("Take the Online AI-900 Practice Exams", "Admin")
        self.ref_head(slide, "Take the Online AI-900 Practice Exams", "Practice Exam", BLUE)
        items = [
            "Reinforce AI-900 topics with realistic, exam-style questions before assessment day.",
            "Use practice sets to test AI workloads, responsible AI, machine learning, Azure AI services and generative AI concepts.",
            "Review mode helps you check the correct answer and revisit the concept after each question.",
            "Complete the practice questions from the Tertiary Exams portal, then bring difficult items to the review session.",
        ]
        self.bullets(slide, Inches(0.92), Inches(2.42), Inches(4.95), Inches(3.15), items, size=15.5, gap=13)
        self.text(slide, Inches(0.92), Inches(5.86), Inches(5.0), Inches(0.28), [[("Access and practise at:", 15, INK, True)]])
        self.text(slide, Inches(0.92), Inches(6.18), Inches(5.45), Inches(0.32), [[("exams.tertiaryinfotech.com", 17, BLUE, True)]])
        self.rrect(slide, Inches(6.10), Inches(2.36), Inches(6.45), Inches(4.70), WHITE, line=RGBColor(0xC9, 0xD6, 0xE6))
        if TERTIARY_MARK.exists():
            slide.shapes.add_picture(str(TERTIARY_MARK), Inches(6.28), Inches(2.51), height=Inches(0.24))
        self.text(slide, Inches(6.58), Inches(2.51), Inches(2.0), Inches(0.18), [[("Tertiary Exams", 7.2, INK, True)]])
        self.text(slide, Inches(8.62), Inches(2.51), Inches(2.7), Inches(0.18), [[("Practice Exam           About Us              Sign In", 5.3, INK, False)]])
        self.rrect(slide, Inches(11.68), Inches(2.48), Inches(0.58), Inches(0.24), BLUE, line=None)
        self.text(slide, Inches(11.70), Inches(2.53), Inches(0.54), Inches(0.09), [[("Get started", 4.8, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(6.43), Inches(3.20), Inches(4.9), Inches(0.14), [[("All exams / Microsoft / AI-900", 5.2, BLUE, True)]])
        self.rrect(slide, Inches(6.43), Inches(3.36), Inches(0.46), Inches(0.18), RGBColor(0xE8, 0xF2, 0xFF), line=RGBColor(0xBF, 0xD7, 0xFF))
        self.text(slide, Inches(6.48), Inches(3.40), Inches(0.36), Inches(0.07), [[("AI-900", 3.8, BLUE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(6.43), Inches(3.62), Inches(4.3), Inches(0.28), [[(TITLE_WITH_PARENS + " Practice Exams", 12, INK, True)]])
        self.text(slide, Inches(6.43), Inches(3.98), Inches(4.3), Inches(0.34), [[("Practice questions for AI workloads, responsible AI, machine learning, Azure AI services and generative AI review.", 5.4, GREY, False)]])
        cards = [("MODE", "Practice", BLUE), ("FOCUS", "AI-900", TEAL), ("ACCESS", "Exam portal", RGBColor(0x7C, 0x3A, 0xED))]
        for idx, (h, v, col) in enumerate(cards):
            x = Inches(6.43 + idx * 1.28)
            self.rrect(slide, x, Inches(4.54), Inches(1.15), Inches(0.54), WHITE, line=LINE)
            self.text(slide, x + Inches(0.12), Inches(4.68), Inches(0.9), Inches(0.11), [[(h, 5.6, GREY, True)]])
            self.text(slide, x + Inches(0.12), Inches(4.88), Inches(0.9), Inches(0.11), [[(v, 7.0, INK, True)]])
            self.rect(slide, x + Inches(0.12), Inches(5.04), Inches(0.85), Inches(0.03), col)
        self.rrect(slide, Inches(6.43), Inches(5.62), Inches(3.95), Inches(0.90), WHITE, line=LINE)
        self.text(slide, Inches(6.58), Inches(5.77), Inches(3.5), Inches(0.16), [[("Domains covered", 7.5, INK, True)]])
        domains = [("AI workloads", 0.90, BLUE), ("Responsible AI", 0.84, TEAL), ("Azure AI services", 0.78, RGBColor(0x7C, 0x3A, 0xED)), ("Generative AI", 0.72, RGBColor(0xF5, 0x9E, 0x0B))]
        yy = 6.05
        for name, pct, col in domains:
            self.text(slide, Inches(6.58), Inches(yy), Inches(1.45), Inches(0.08), [[(name, 4.8, INK, False)]])
            self.rect(slide, Inches(8.02), Inches(yy + 0.025), Inches(1.55), Inches(0.04), RGBColor(0xD9, 0xE1, 0xEA))
            self.rect(slide, Inches(8.02), Inches(yy + 0.025), Inches(1.55 * pct), Inches(0.04), col)
            yy += 0.14
        self.rrect(slide, Inches(10.73), Inches(3.58), Inches(1.28), Inches(0.94), WHITE, line=LINE)
        self.text(slide, Inches(10.90), Inches(3.83), Inches(0.95), Inches(0.13), [[("PRACTICE DETAILS", 5.2, INK, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(10.90), Inches(4.08), Inches(0.95), Inches(0.13), [[("Portal: exams.tertiaryinfotech.com", 4.4, GREY, False)]], align=PP_ALIGN.CENTER)
        self.rrect(slide, Inches(10.73), Inches(5.36), Inches(1.28), Inches(0.94), RGBColor(0xE8, 0xF2, 0xFF), line=RGBColor(0xBF, 0xD7, 0xFF))
        self.text(slide, Inches(10.90), Inches(5.61), Inches(0.95), Inches(0.15), [[("Learner tip", 6.0, BLUE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(10.90), Inches(5.86), Inches(0.95), Inches(0.15), [[("Retake missed items after revision.", 4.8, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(slide)

    def course_repository_slide(self):
        slide = self._slide("Access the Hands-On Labs", "Admin")
        self.head(slide, "Access the Hands-On Labs", "COURSE REPOSITORY", BLUE)
        self.text(slide, Inches(0.9), Inches(1.62), Inches(11.4), Inches(0.48), [[("Use the GitHub repository to review lab files, revisit steps after class, and keep your Learner Guide aligned with the exercises.", 14, GREY, False)]])
        self.rrect(slide, Inches(0.9), Inches(2.35), Inches(11.45), Inches(0.85), LIGHT, line=TEAL)
        self.text(slide, Inches(1.15), Inches(2.62), Inches(11.0), Inches(0.3), [[(f"https://{LAB_REPO}", 14, BLUE, True)]])
        opts = [
            ("Option A", "Clone the repository with git clone <course-repo-url>.git"),
            ("Option B", "Download a ZIP from GitHub: Code > Download ZIP"),
            ("Run in class", "Use Azure portal, Microsoft Learn, and instructor-provided accounts as directed."),
        ]
        for idx, (head, body) in enumerate(opts):
            x = Inches(0.9) + idx * Inches(3.85)
            self.rrect(slide, x, Inches(3.75), Inches(3.55), Inches(1.45), WHITE, line=LINE)
            self.text(slide, x + Inches(0.18), Inches(3.98), Inches(3.15), Inches(0.28), [[(head, 13, TEAL, True)]])
            self.text(slide, x + Inches(0.18), Inches(4.38), Inches(3.15), Inches(0.55), [[(body, 11.5, INK, False)]])
        self.footer(slide)

    def lab_overview(self, lab: Lab):
        slide = self._slide(f"Lab {lab.num}: {lab.title}", f"Lab {lab.num}")
        self.head(slide, lab.title, f"LAB {lab.num}", TEAL)
        self.rect(slide, Inches(0.9), Inches(1.75), Inches(1.45), Inches(0.45), TEAL)
        self.text(slide, Inches(0.9), Inches(1.82), Inches(1.45), Inches(0.3), [[(f"LAB {lab.num:02d}", 14, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.text(slide, Inches(2.6), Inches(1.78), Inches(9.8), Inches(0.5), [[(lab.scenario[:220], 14.5, GREY, False)]])
        self.text(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.35), [[("Objectives", 16, BLUE, True)]])
        self.bullets(slide, Inches(1.05), Inches(3.0), Inches(11.1), Inches(2.15), lab.objectives[:5], size=15, gap=6)
        self.rrect(slide, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.95), LIGHT)
        self.text(slide, Inches(1.15), Inches(5.55), Inches(11.0), Inches(0.45), [[("Validation: ", 14, TEAL, True), (lab.validation[:210], 14, INK, False)]])
        self.footer(slide)

    def lab_steps(self, lab: Lab):
        slide = self._slide(f"Lab {lab.num} guided steps", f"Lab {lab.num}")
        self.head(slide, f"Lab {lab.num} guided steps", f"LAB {lab.num} - STEP FLOW", TEAL)
        items = lab.step_titles[:7]
        if len(lab.step_titles) > 7:
            items.append("Complete validation, checkpoints, and cleanup review.")
        self.bullets(slide, Inches(0.95), Inches(1.75), Inches(11.4), Inches(4.75), items, size=17, marker=TEAL)
        self.footer(slide)

    def lab_flow_diagram(self, lab: Lab):
        slide = self._slide(f"Lab {lab.num} workflow diagram", f"Lab {lab.num}")
        self.head(slide, f"Lab {lab.num} workflow diagram", f"LAB {lab.num} - DIAGRAM", TEAL)
        steps = lab.step_titles[:5] or ["Interpret scenario", "Complete guided tasks", "Validate output", "Review controls"]
        if len(lab.step_titles) > 5:
            steps[-1] = "Validate and review"
        count = len(steps)
        gap = Inches(0.28)
        box_w = (Inches(11.5) - gap * (count - 1)) // count
        y = Inches(2.35)
        for idx, step in enumerate(steps):
            x = Inches(0.9) + idx * (box_w + gap)
            self.rrect(slide, x, y, box_w, Inches(1.55), WHITE, line=TEAL if idx % 2 else BLUE)
            self.oval(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.42), Inches(0.42), TEAL if idx % 2 else BLUE)
            self.text(slide, x + Inches(0.18), y + Inches(0.24), Inches(0.42), Inches(0.22), [[(str(idx + 1), 11, WHITE, True)]], align=PP_ALIGN.CENTER)
            self.text(slide, x + Inches(0.18), y + Inches(0.73), box_w - Inches(0.36), Inches(0.62), [[(step[:48], 11.5, INK, True)]], align=PP_ALIGN.CENTER)
            if idx < count - 1:
                self.text(slide, x + box_w + Inches(0.03), y + Inches(0.52), gap - Inches(0.05), Inches(0.3), [[(">", 20, BLUE, True)]], align=PP_ALIGN.CENTER)
        self.rrect(slide, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.15), LIGHT)
        self.text(slide, Inches(1.2), Inches(4.78), Inches(10.9), Inches(0.55), [[("Lab evidence: ", 15, BLUE, True), (lab.validation[:240], 15, INK, False)]])
        self.footer(slide)

    def lab_review(self, lab: Lab):
        slide = self._slide(f"Lab {lab.num} checkpoint", f"Lab {lab.num}")
        self.head(slide, f"Lab {lab.num} checkpoint", f"LAB {lab.num} - REVIEW", TEAL)
        left = lab.checkpoint[:4] or ["Confirm the lab artifact is complete.", "Explain the responsible AI implication."]
        right = [lab.focus] if lab.focus else ["Review the workload, the service, and the human review point."]
        self.rrect(slide, Inches(0.85), Inches(1.85), Inches(5.65), Inches(4.65), LIGHT)
        self.rrect(slide, Inches(6.85), Inches(1.85), Inches(5.65), Inches(4.65), LIGHT)
        self.text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.35), [[("Checkpoint questions", 16, BLUE, True)]])
        self.bullets(slide, Inches(1.1), Inches(2.65), Inches(5.0), Inches(3.5), left, size=14)
        self.text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.35), [[("Course focus", 16, TEAL, True)]])
        self.text(slide, Inches(7.1), Inches(2.7), Inches(5.0), Inches(2.2), [[(right[0][:330], 16, INK, False)]])
        self.footer(slide)

    def save(self, path: Path) -> None:
        self.prs.save(path)


def build_slides(labs: list[Lab]) -> tuple[Path, list[dict[str, str | int]]]:
    b = SlideBuilder()
    b.cover()
    b.admin_slide("Welcome & Housekeeping", "Course Administration", [
        "Confirm attendance, facilities, emergency procedures, and course timing.",
        "Use the LMS/TMS portal for courseware, digital attendance, feedback, and assessment.",
        "Keep this deck and Learner Guide available during class activities.",
    ])
    b.admin_slide("Digital Attendance (Mandatory)", "TRAQOM - SSG Digital Attendance", [
        "Take AM, PM, and assessment digital attendance when the QR code is displayed.",
        "Scan the QR code with your mobile phone camera and submit the attendance form.",
        "Minimum 75% attendance is required for funding eligibility.",
    ])
    b.trainer_template_slide()
    b.trainer_profile_slide()
    b.admin_slide("Let's Know Each Other", "Ice-breaker", [
        "Your name and organisation / role.",
        "Your experience with Azure, AI, analytics, or cloud services.",
        "One AI scenario you want to understand by the end of the course.",
    ])
    b.ground_rules_slide()
    b.admin_slide("LMS / TMS", "Course Portal", [
        "Portal: https://lms-tms.tertiaryinfotech.com",
        "Download slides and Learner Guide from the portal.",
        "Courseware and approved materials may be used for the open-book assessment.",
    ])
    b.admin_slide("Download the Labs", "GitHub Course Repository", [
        f"Open: {LAB_REPO}",
        "Complete the labs in order because later activities reuse the same Azure AI vocabulary.",
        "Use an instructor-provided Azure subscription, Microsoft Learn sandbox, or Azure free account.",
    ])
    b.two_col("Lesson Plan", "Day 1", "Morning", [
        "Digital attendance (AM)",
        "Trainer and learner introductions",
        "AI workloads and responsible AI",
        "Lab 01: AI workloads and responsible AI",
        "Lab 02: Machine learning fundamentals",
    ], "Afternoon", [
        "Digital attendance (PM)",
        "Azure Machine Learning and model lifecycle",
        "Lab 03: Azure Machine Learning",
        "Lab 04: Computer vision, image analysis, OCR",
        "Day 1 review",
    ])
    b.two_col("Lesson Plan", "Day 2", "Morning", [
        "Digital attendance (AM)",
        "NLP, speech, translation, document intelligence",
        "Labs 05-07: Azure AI services",
        "Generative AI, Foundry, OpenAI, model catalog",
    ], "Afternoon", [
        "Lab 08: Generative AI",
        "Lab 09: Solution design and governance",
        "Lab 10: Capstone and course review",
        "Digital attendance (Assessment) and final assessment",
        "TRAQOM survey and closing",
    ])
    b.two_col("Skills Framework", "WSQ TSC", "TSC Title", [
        f"{TSC_TITLE}",
        f"TSC Code: {TSC_CODE}",
        "Cloud and AI service awareness for business and technical roles",
        "Responsible use of AI-enabled services",
    ], "Applied Abilities", [
        "Identify AI workloads from business requirements.",
        "Select suitable Azure AI services.",
        "Recognize governance, privacy, safety, and cost controls.",
    ])
    b.tsc_detail_slide("TSC Knowledge (K)", "Skills Framework", [
        ("K1", "Artificial intelligence workload types and business use cases."),
        ("K2", "Responsible AI principles, privacy, security, transparency, and accountability."),
        ("K3", "Machine learning concepts including features, labels, training, evaluation, and deployment."),
        ("K4", "Azure AI service capabilities for vision, language, speech, documents, and search."),
        ("K5", "Generative AI concepts including prompts, foundation models, grounding, and content safety."),
        ("K6", "Governance, access control, monitoring, cost, and resource cleanup considerations."),
    ], f"TSC Title: {TSC_TITLE} | TSC Code: {TSC_CODE}")
    b.tsc_detail_slide("TSC Abilities (A)", "Skills Framework", [
        ("A1", "Identify AI workloads from practical business scenarios and input/output requirements."),
        ("A2", "Map Azure AI services to workload requirements and learner lab evidence."),
        ("A3", "Apply responsible AI review questions to vision, language, speech, document, and generative AI use cases."),
        ("A4", "Explain machine learning and generative AI concepts using simple workplace examples."),
        ("A5", "Recommend basic security, privacy, monitoring, and cost controls for AI solutions."),
        ("A6", "Prepare a capstone service map and review plan aligned to the labs."),
    ], "Abilities are practiced through Labs 01-10 and referenced in the Lesson Plan slide numbers.")
    b.skills_framework_diagram()
    b.content_slide("Learning Outcomes", "Course Outcomes", [
        "Identify common AI workloads and use cases.",
        "Explain responsible AI principles and human review needs.",
        "Distinguish regression, classification, clustering, deep learning, and transformer-based workloads.",
        "Describe Azure Machine Learning and automated ML concepts.",
        "Map vision, NLP, speech, translation, document, and generative AI scenarios to Azure services.",
        "Recognize security, privacy, governance, and cost controls for AI solutions.",
    ], size=16)
    b.two_col("Course Outline", "Topics 1-4", "Foundations", [
        "AI workloads and responsible AI",
        "Machine learning fundamentals",
        "Features, labels, training, evaluation, deployment",
        "Azure Machine Learning and AutoML",
    ], "Vision and Language", [
        "Image analysis, OCR, object detection",
        "Sentiment, key phrase extraction, entity recognition",
        "Speech recognition, speech synthesis, translation",
    ])
    b.two_col("Course Outline", "Topics 5-8", "Documents and Generative AI", [
        "Document Intelligence and knowledge mining",
        "Azure AI Search concepts",
        "Generative AI, prompts, LLMs, grounding, RAG",
        "Azure AI Foundry, Azure OpenAI, model catalog",
    ], "Governance and Review", [
        "Security, privacy, managed identity, RBAC",
        "Budget alerts and resource cleanup",
        "Capstone service mapping",
        "Assessment preparation and course review",
    ])
    b.admin_slide("Final Assessment", "Assessment", [
        "Format: written / MCQ assessment on the LMS.",
        "Assessment is open book: approved course slides, Learner Guide, and materials only.",
        "Complete the assessment digital attendance before starting.",
        "Appeals follow the centre's published assessment process.",
    ])
    b.admin_slide("Briefing for Assessment", "Assessment Rules", [
        "Place phones and other materials under the table or as instructed.",
        "No photos or recording of assessment scripts.",
        "No discussion during assessment.",
        "Submit answers on the LMS before the time ends.",
        "Sign the Assessment Summary Record when instructed.",
    ])
    b.admin_slide("Criteria for Funding", "Funding", [
        "Minimum attendance rate of 75% based on SSG digital attendance records.",
        "Complete the assessment and be assessed as Competent.",
        "Complete course feedback and TRAQOM survey when requested.",
    ])
    b.course_repository_slide()
    b.admin_slide("Certification & TRAQOM Survey", "Before You Go", [
        "Complete the certification and TRAQOM survey to receive your certificate.",
        "Confirm AM, PM, and assessment attendance are recorded.",
        "Keep the Learner Guide for post-course revision and workplace reference.",
    ])

    b.section("Course Concepts", "AI-900 builds recognition skill through scenarios", "Learners practice mapping business needs to AI workloads, Azure services, and responsible controls.")
    b.content_slide("Course Status Note", "Important", [
        "Microsoft Learn states that Exam AI-900 retired on June 30, 2026.",
        "This courseware remains aligned to the AI-900 skills outline and Azure AI fundamentals concepts.",
        "The course should be positioned as Azure AI Fundamentals training, not as a guarantee of exam appointment availability.",
    ], size=18, accent=AMBER)
    b.content_slide("AI workload recognition is the core skill", "Concept", [
        "Prediction and classification: use data to estimate values or categories.",
        "Vision: interpret images, detect objects, and read text with OCR.",
        "Language and speech: understand text, translate, transcribe, and synthesize voice.",
        "Documents: extract structured information and make unstructured content searchable.",
        "Generative AI: draft, summarize, answer, and assist using foundation models.",
    ], size=17)
    b.diagram_slide()
    b.content_slide("Responsible AI keeps automation accountable", "Concept", [
        "Fairness: evaluate whether different groups are affected unequally.",
        "Reliability and safety: test outputs and define fallback paths.",
        "Privacy and security: protect data, access, keys, and logs.",
        "Inclusiveness: design for different languages, abilities, and contexts.",
        "Transparency and accountability: explain when AI is used and who owns decisions.",
    ], size=17)
    b.content_slide("Machine learning follows a repeatable lifecycle", "Concept", [
        "Define the prediction task and business success measure.",
        "Prepare features and labels from trustworthy training data.",
        "Train and evaluate a model using appropriate metrics.",
        "Deploy, monitor, retrain, and govern the model over time.",
    ], size=18)
    b.content_slide("Generative AI adds new controls", "Concept", [
        "Prompts guide the model, but outputs still need review.",
        "Grounding and retrieval reduce unsupported answers by adding source context.",
        "Content filters, evaluations, and monitoring help manage harmful or inaccurate outputs.",
        "Users should know when content is AI-assisted.",
    ], size=18, accent=TEAL)

    for lab in labs:
        b.lab_overview(lab)
        b.lab_flow_diagram(lab)
        b.lab_steps(lab)
        b.lab_review(lab)

    b.section("Course Synthesis", "The capstone connects workload, service, and control", "A useful Azure AI recommendation names the workload, the service, the data boundary, and the responsible AI guardrail.")
    b.content_slide("Final service selection checklist", "Review", [
        "What is the input: image, text, speech, document, tabular data, or prompt?",
        "What output is needed: prediction, label, extracted field, transcript, translation, summary, or answer?",
        "Which Azure service is the closest fit?",
        "What privacy, security, human review, and cost controls are required?",
        "How will the learner validate the result and clean up resources?",
    ], size=17)
    b.assessment_flow()
    b.practice_exam_slide()
    b.content_slide("Course close", "Before You Go", [
        "Complete TRAQOM survey and final attendance.",
        "Save the Learner Guide and lab notes for post-course review.",
        "Clean up personal Azure resources or confirm shared cleanup with the trainer.",
        "Review weak areas using the 7-day study plan in Lab 10.",
    ], size=18)

    out = COURSEWARE / f"{TITLE}-{VERSION}.pptx"
    try:
        b.save(out)
    except PermissionError:
        out = COURSEWARE / f"{TITLE}-{VERSION}-updated.pptx"
        b.save(out)
    return out, b.slide_index


def set_cell_bg(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_doc_run(paragraph, text: str) -> None:
    for part in re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            r = paragraph.add_run(part[2:-2])
            r.bold = True
        elif part.startswith("`") and part.endswith("`"):
            r = paragraph.add_run(part[1:-1])
            r.font.name = "Consolas"
            r.font.color.rgb = DRGBColor(0xC7, 0x25, 0x4E)
        else:
            paragraph.add_run(part)


def style_doc(doc: Document) -> None:
    sec = doc.sections[0]
    sec.page_width = DInches(8.27)
    sec.page_height = DInches(11.69)
    sec.top_margin = DInches(0.65)
    sec.bottom_margin = DInches(0.65)
    sec.left_margin = DInches(0.7)
    sec.right_margin = DInches(0.7)
    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)
    for name, size, color in [
        ("Heading 1", 16, DRGBColor(0x1F, 0x6F, 0xEB)),
        ("Heading 2", 13, DRGBColor(0x11, 0x18, 0x27)),
        ("Heading 3", 11.5, DRGBColor(0x1F, 0x6F, 0xEB)),
    ]:
        st = styles[name]
        st.font.name = "Arial"
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = color


def add_field(paragraph, instr: str, default: str = "") -> None:
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    it = OxmlElement("w:instrText")
    it.set(qn("xml:space"), "preserve")
    it.text = instr
    sep = OxmlElement("w:fldChar")
    sep.set(qn("w:fldCharType"), "separate")
    txt = OxmlElement("w:t")
    txt.text = default
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(begin)
    run._r.append(it)
    run._r.append(sep)
    run._r.append(txt)
    run._r.append(end)


def cover(doc: Document, kind: str) -> None:
    title_text = kind.upper()
    if TERTIARY_MARK.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(38)
        p.paragraph_format.space_after = Pt(4)
        p.add_run().add_picture(str(TERTIARY_MARK), width=DInches(1.10))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(0)
    r = p.add_run(ORG)
    r.bold = True
    r.font.size = Pt(10)
    r.font.color.rgb = DRGBColor(0x11, 0x18, 0x27)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(10)
    r = p.add_run("UEN: 201200696W")
    r.font.size = Pt(8)
    r.font.color.rgb = DRGBColor(0x55, 0x5B, 0x66)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(title_text)
    r.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = DRGBColor(0x1F, 0x6F, 0xEB)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(10)
    r = p.add_run("For")
    r.font.size = Pt(8)
    r.font.color.rgb = DRGBColor(0x55, 0x5B, 0x66)
    if COURSE_BADGE.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(14)
        p.add_run().add_picture(str(COURSE_BADGE), width=DInches(0.95))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(TITLE_WITH_PARENS)
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = DRGBColor(0x11, 0x18, 0x27)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(8)
    r = p.add_run(f"TGS Ref No: {CODE}")
    r.font.size = Pt(9)
    if "LEARNER" in title_text:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(0)
        r = p.add_run("Conducted by")
        r.font.size = Pt(8)
        r.font.color.rgb = DRGBColor(0x55, 0x5B, 0x66)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(0)
        r = p.add_run(ORG)
        r.bold = True
        r.font.size = Pt(10)
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(8)
        r = p.add_run("UEN: 201200696W")
        r.font.size = Pt(8)
        r.font.color.rgb = DRGBColor(0x55, 0x5B, 0x66)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"Version {DOC_VERSION} - {DATE}")
    r.bold = True
    r.font.size = Pt(10)
    r.font.color.rgb = DRGBColor(0x1F, 0x6F, 0xEB)
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def version_control(doc: Document, summary: str) -> None:
    p = doc.add_paragraph()
    r = p.add_run("DOCUMENT VERSION CONTROL RECORD")
    r.bold = True
    r.font.size = Pt(12)
    t = doc.add_table(rows=1, cols=4)
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["Version Number", "Effective Date of Release", "Summary of Included Changes", "Author"]
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        c.text = h
        set_cell_bg(c, "1F6FEB")
        for run in c.paragraphs[0].runs:
            run.bold = True
            run.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(9.5)
    row = t.add_row().cells
    for c, val in zip(row, [DOC_VERSION, DATE, summary, ORG]):
        c.text = val
        for run in c.paragraphs[0].runs:
            run.font.size = Pt(9.5)
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def toc(doc: Document, levels: str = "1-3") -> None:
    p = doc.add_paragraph()
    r = p.add_run("TABLE OF CONTENTS")
    r.bold = True
    r.font.size = Pt(12)
    p = doc.add_paragraph()
    add_field(p, f'TOC \\o "{levels}" \\h \\z \\u', "Right-click and choose Update Field to build the table of contents.")
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def footer(doc: Document) -> None:
    sec = doc.sections[0]
    p = sec.footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.text = ""
    p.add_run("Page ")
    add_field(p, "PAGE", "1")
    p.add_run(" of ")
    add_field(p, "NUMPAGES", "1")
    p2 = sec.footer.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p2.add_run(f"(c) 2026 {ORG}. All rights reserved. - {URL}")
    r.font.size = Pt(7.5)
    r.font.color.rgb = DRGBColor(0x55, 0x5B, 0x66)


def render_markdown_docx(doc: Document, md_lines: list[str]) -> None:
    in_code = False
    code_buf: list[str] = []
    for raw in md_lines:
        line = raw.rstrip()
        if line.startswith("```"):
            if in_code:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = DInches(0.2)
                p.paragraph_format.space_before = Pt(3)
                p.paragraph_format.space_after = Pt(6)
                r = p.add_run("\n".join(code_buf))
                r.font.name = "Consolas"
                r.font.size = Pt(8.5)
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue
        if not line.strip():
            continue
        img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
        if img_match:
            img_path = COURSEWARE / img_match.group(2)
            if img_path.exists():
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(8)
                p.add_run().add_picture(str(img_path), width=DInches(6.6))
            continue
        if line.startswith("### "):
            doc.add_paragraph(strip_md(line[4:]), style="Heading 3")
        elif line.startswith("## "):
            doc.add_paragraph(strip_md(line[3:]), style="Heading 2")
        elif line.startswith("# "):
            doc.add_paragraph(strip_md(line[2:]), style="Heading 1")
        elif line.startswith("- "):
            p = doc.add_paragraph(style="List Bullet")
            add_doc_run(p, line[2:])
        elif re.match(r"^\d+\.\s+", line):
            p = doc.add_paragraph(style="List Number")
            add_doc_run(p, re.sub(r"^\d+\.\s+", "", line))
        elif line.startswith("| ") and line.endswith(" |"):
            # Tables are retained in the Markdown mirror; in DOCX they are rendered as compact monospaced rows.
            p = doc.add_paragraph()
            r = p.add_run(strip_md(line))
            r.font.name = "Consolas"
            r.font.size = Pt(8)
        else:
            p = doc.add_paragraph()
            add_doc_run(p, line)


def build_lg(labs: list[Lab]) -> tuple[Path, Path]:
    ensure_lg_diagrams(labs)
    md_path = COURSEWARE / f"LG-{TITLE}.md"
    docx_path = COURSEWARE / f"LG-{TITLE}.docx"
    summary = "First AI-900 learner guide aligned to the 10 GitHub labs; includes step-by-step lab guidance and Markdown mirror."
    chunks: list[str] = []
    chunks.append(f"# Learner Guide - {TITLE_WITH_PARENS}\n")
    chunks.append(f"**Course Code:** {CODE}  \n**Version:** {DOC_VERSION}  \n**Provider:** {ORG}\n")
    chunks.append("## Course Information\n")
    chunks.append(f"| Item | Details |\n| --- | --- |\n| Course Title | {TITLE_WITH_PARENS} |\n| Course Code | {CODE} |\n| WSQ TSC | {TSC_TITLE} ({TSC_CODE}) |\n| Duration | 2 Days (16 training hours) |\n| Delivery Mode | Instructor-led with guided Azure AI lab activities |\n| Course Repository | https://{LAB_REPO} |\n")
    chunks.append("## Certification Status Note\n")
    chunks.append("Microsoft Learn states that Exam AI-900 was retired on June 30, 2026. These materials are therefore written as Azure AI Fundamentals courseware aligned to the AI-900 skills outline and Azure AI concepts.\n")
    chunks.append("## Course Goal\n")
    chunks.append("This course helps learners recognize common AI workloads, explain responsible AI considerations, and map business scenarios to Azure AI services through hands-on lab activities.\n")
    chunks.append("## Course Learning Flow\n")
    chunks.append("![Course learning flow](assets/lg-diagrams/course-learning-flow.png)\n")
    chunks.append("## Learning Outcomes\n")
    chunks.append("- Identify common AI workloads and use cases.\n- Explain responsible AI principles and human review needs.\n- Distinguish regression, classification, clustering, deep learning, and transformer-based workloads.\n- Describe Azure Machine Learning and automated ML concepts.\n- Map vision, NLP, speech, translation, document, and generative AI scenarios to Azure services.\n- Recognize security, privacy, governance, and cost controls for AI solutions.\n")
    chunks.append("## Skills Framework Alignment\n")
    chunks.append(f"**TSC Title:** {TSC_TITLE}  \n**TSC Code:** {TSC_CODE}\n")
    chunks.append("| Knowledge / Ability | Lab Alignment |\n| --- | --- |\n| AI workload recognition and responsible AI | Labs 01, 09, 10 |\n| Machine learning concepts and Azure Machine Learning | Labs 02, 03 |\n| Vision, language, speech, and document AI services | Labs 04, 05, 06, 07 |\n| Generative AI, Azure AI Foundry, Azure OpenAI, and model catalog concepts | Lab 08 |\n| Governance, security, privacy, cost, and cleanup controls | Labs 09, 10 |\n")
    chunks.append("## Azure AI Service Map\n")
    chunks.append("![Azure AI service map](assets/lg-diagrams/azure-ai-service-map.png)\n")
    chunks.append("## Lab Environment Setup\n")
    chunks.append("1. Open https://portal.azure.com/.\n2. Sign in with an instructor-provided or personal Azure account.\n3. Confirm subscription and resource group access.\n4. Use free tiers where possible.\n5. Do not paste personal or confidential data into AI demos.\n6. Clean up resources when the trainer confirms it is safe.\n")
    chunks.append("## Labs\n")
    for lab in labs:
        source = (LABS_DIR / lab.filename).read_text(encoding="utf-8")
        source = re.sub(r"^# Lab", "## Lab", source, count=1, flags=re.M)
        source = source.strip() + f"\n\n### Lab Workflow Diagram\n\n![Lab {lab.num:02d} workflow diagram](assets/lg-diagrams/lab-{lab.num:02d}-workflow.png)\n"
        chunks.append(source + "\n")
    chunks.append("## Final Review Checklist\n")
    chunks.append("- Match AI workloads to Azure services.\n- Explain responsible AI principles.\n- Distinguish machine learning techniques.\n- Recognize vision, language, speech, document, and generative AI scenarios.\n- Identify security, privacy, cost, and governance considerations.\n")
    md = "\n".join(chunks)
    md_path.write_text(md, encoding="utf-8")

    doc = Document()
    style_doc(doc)
    cover(doc, "Learner Guide")
    version_control(doc, summary)
    toc(doc)
    render_markdown_docx(doc, md.splitlines())
    footer(doc)
    doc.settings.element.append(OxmlElement("w:updateFields"))
    doc.save(docx_path)
    return md_path, docx_path


def find_slide(slide_index: list[dict[str, str | int]], contains: str) -> int:
    for row in slide_index:
        if contains.lower() in str(row["title"]).lower():
            return int(row["slide"])
    return 1


def lab_slide_range(slide_index: list[dict[str, str | int]], lab_num: int) -> str:
    nums = [int(r["slide"]) for r in slide_index if str(r["section"]) == f"Lab {lab_num}"]
    if not nums:
        return "-"
    return f"{min(nums)}-{max(nums)}"


def lp_day_rows(slide_index: list[dict[str, str | int]]) -> dict[str, list[tuple[str, str, str, str, str]]]:
    admin = f"1-{find_slide(slide_index, 'Certification & TRAQOM Survey')}"
    concept = f"{find_slide(slide_index, 'AI-900 builds')}-{find_slide(slide_index, 'Generative AI adds')}"
    assess = f"{find_slide(slide_index, 'Final Assessment')}, {find_slide(slide_index, 'Briefing for Assessment')}, {find_slide(slide_index, 'Assessment Flow')}-{find_slide(slide_index, 'Take the Online AI-900 Practice Exams')}"
    close = str(find_slide(slide_index, "Course close"))
    return {
        "Day 1 - AI and Machine Learning Foundations": [
            ("9:00-9:15", "Course opening and attendance", "Confirm TRAQOM attendance, class logistics, assessment rules, courseware access, and learner expectations.", "Learners scan attendance, confirm access to slides/LG/labs, and state role or AI experience.", admin),
            ("9:15-9:45", "AI workloads and responsible AI", "Explain workload categories with workplace examples; prompt learners to classify each scenario by input, output, and risk.", "Learners classify scenarios and record responsible AI questions for later lab checkpoints.", concept),
            ("9:45-10:30", "Lab 01: AI workloads and responsible AI", "Guide learners through workload recognition and responsible AI review; pause after each checkpoint for discussion.", "Completed Lab 01 notes: workload selected, responsible AI issue identified, and validation response prepared.", lab_slide_range(slide_index, 1)),
            ("10:30-10:45", "Tea break", "Remind learners to save lab notes and keep Azure/Microsoft Learn sessions active.", "Break.", "-"),
            ("10:45-11:30", "Lab 02: Machine learning fundamentals", "Demonstrate features, labels, training data, evaluation, and prediction vocabulary before hands-on practice.", "Learners explain the difference between classification, regression, clustering, and model evaluation.", lab_slide_range(slide_index, 2)),
            ("11:30-12:15", "Lab 03: Azure Machine Learning", "Walk through Azure Machine Learning/AutoML concepts and connect portal steps to the ML lifecycle.", "Learners identify model lifecycle stages and note where evaluation and deployment evidence appears.", lab_slide_range(slide_index, 3)),
            ("12:15-1:00", "Guided review: ML lifecycle and service mapping", "Use the services map to compare Azure Machine Learning with prebuilt Azure AI services.", "Learners complete a mini service-selection exercise and justify the service chosen.", concept),
            ("1:00-2:00", "Lunch break", "Confirm afternoon lab sequence and support learners who need account access help.", "Break.", "-"),
            ("2:00-2:45", "Lab 04: Computer vision, image analysis, OCR", "Brief image inputs, object/text extraction, confidence scores, and human review before the lab.", "Learners validate image/OCR output and list accuracy or privacy risks.", lab_slide_range(slide_index, 4)),
            ("2:45-3:30", "Lab 05: Natural language processing", "Explain sentiment, key phrase extraction, entity recognition, and language service boundaries.", "Learners compare NLP output with the source text and capture a review question.", lab_slide_range(slide_index, 5)),
            ("3:30-3:45", "Tea break", "Check progress and identify learners needing remediation before the workshop.", "Break.", "-"),
            ("3:45-4:30", "Scenario workshop", "Facilitate small-group service selection across image, text, and ML scenarios.", "Learners produce a workload-service-control mapping for at least three scenarios.", f"{concept}, {lab_slide_range(slide_index, 4)}-{lab_slide_range(slide_index, 5)}"),
            ("4:30-5:15", "Checkpoint and remediation", "Review wrong answers, address misunderstood vocabulary, and connect answers back to lab evidence.", "Learners update lab notes and complete checkpoint questions.", "Lab review slides"),
            ("5:15-6:00", "Day 1 recap and Q&A", "Summarize workload recognition, ML lifecycle, vision, and language services; assign revision topics.", "Learners state one concept mastered and one concept to review on Day 2.", "Review slides"),
        ],
        "Day 2 - Azure AI Services, Generative AI, and Assessment": [
            ("9:00-9:15", "Attendance and Day 1 recap", "Confirm AM attendance and run a quick recap quiz on workload-service mapping.", "Learners answer recap questions and reopen their lab repository/courseware.", admin),
            ("9:15-10:00", "Lab 06: Speech, translation, conversational AI", "Explain speech-to-text, text-to-speech, translation, and conversation use cases before lab work.", "Learners validate transcript/translation outputs and record accessibility or privacy considerations.", lab_slide_range(slide_index, 6)),
            ("10:00-10:45", "Lab 07: Document intelligence and search", "Demonstrate document extraction, knowledge mining, indexing, and search concepts.", "Learners identify extracted fields, indexing value, and when human review is required.", lab_slide_range(slide_index, 7)),
            ("10:45-11:00", "Tea break", "Check lab progress and prepare generative AI safety framing.", "Break.", "-"),
            ("11:00-12:00", "Lab 08: Generative AI, Foundry, Azure OpenAI", "Introduce prompts, grounding, content safety, model catalog, and review controls.", "Learners test prompts, evaluate responses, and document grounding/review safeguards.", lab_slide_range(slide_index, 8)),
            ("12:00-1:00", "Guided practice: generative AI controls", "Facilitate comparison of prompt-only, grounded, and reviewed outputs; highlight risk controls.", "Learners refine prompts and write a short governance note for generative AI use.", concept),
            ("1:00-2:00", "Lunch break", "Confirm remaining lab and assessment timing.", "Break.", "-"),
            ("2:00-2:45", "Lab 09: solution design and governance", "Coach learners through service selection, RBAC/privacy/cost controls, and cleanup decisions.", "Learners complete a solution design map and name at least three controls.", lab_slide_range(slide_index, 9)),
            ("2:45-3:30", "Lab 10: capstone and review", "Guide the capstone service map and 7-day revision plan; validate all lab evidence is complete.", "Learners submit/review the capstone map and identify personal revision priorities.", lab_slide_range(slide_index, 10)),
            ("3:30-3:45", "Assessment briefing", "Review allowed materials, digital attendance, LMS/TMS submission, timing, and conduct rules.", "Learners prepare approved materials and complete assessment attendance when instructed.", assess),
            ("3:45-5:00", "Final assessment", "Invigilate the written/MCQ assessment, support LMS/TMS technical access only, and enforce assessment rules.", "Learners complete assessment and submit answers on LMS/TMS.", assess),
            ("5:00-5:15", "Tea break", "Collect any operational issues and prepare closing checks.", "Break.", "-"),
            ("5:15-6:00", "Course close, TRAQOM, practice exam", "Confirm TRAQOM, final attendance, certification actions, resource cleanup, and practice exam access.", "Learners complete required surveys/attendance and save the practice exam link for revision.", f"{find_slide(slide_index, 'Assessment Flow')}-{close}"),
        ],
    }


def lp_lab_alignment_rows(labs: list[Lab], slide_index: list[dict[str, str | int]]) -> list[tuple[str, str, str, str, str]]:
    rows = []
    for lab in labs:
        checkpoint = "; ".join(lab.checkpoint[:2]) if lab.checkpoint else lab.focus or "Confirm validation evidence and explain the service choice."
        focus = lab.focus or lab.scenario
        rows.append((f"{lab.num:02d}", lab.title, lab_slide_range(slide_index, lab.num), focus[:240], f"{lab.validation[:220]} Checkpoint: {checkpoint[:180]}"))
    return rows


def set_table_widths(table, widths: list[float]) -> None:
    for row in table.rows:
        for idx, width in enumerate(widths):
            row.cells[idx].width = DInches(width)


def style_doc_table(table, header_fill: str = "1F6FEB") -> None:
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for cell in table.rows[0].cells:
        set_cell_bg(cell, header_fill)
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
                run.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(8.5)
    for row in table.rows:
        for cell in row.cells:
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_after = Pt(2)
                for run in paragraph.runs:
                    run.font.size = Pt(8.5)


def add_doc_table(doc: Document, headers: list[str], rows: list[tuple[str, ...]], widths: list[float], break_keywords: tuple[str, ...] = ("break",)) -> None:
    tbl = doc.add_table(rows=1, cols=len(headers))
    for i, h in enumerate(headers):
        tbl.rows[0].cells[i].text = h
    for row in rows:
        c = tbl.add_row().cells
        for i, val in enumerate(row):
            c[i].text = str(val)
        if any(k.lower() in " ".join(row).lower() for k in break_keywords):
            for cell in c:
                set_cell_bg(cell, "FFF4E5")
    set_table_widths(tbl, widths)
    style_doc_table(tbl)


def build_lp(labs: list[Lab], slide_index: list[dict[str, str | int]]) -> Path:
    docx_path = COURSEWARE / f"LP-{TITLE}.docx"
    doc = Document()
    style_doc(doc)
    cover(doc, "Lesson Plan")
    version_control(doc, "Detailed AI-900 lesson plan aligned to slide numbers, lab activities, assessment, and delivery checkpoints.")
    toc(doc, "1-2")

    doc.add_paragraph("Course Information", style="Heading 1")
    info = [
        ("Course Title", TITLE_WITH_PARENS),
        ("Course Code", CODE),
        ("WSQ TSC", f"{TSC_TITLE} ({TSC_CODE})"),
        ("Duration", "2 Days (16 training hours)"),
        ("Daily Schedule", "9:00 AM - 6:00 PM (8 training hours/day, excluding lunch)"),
        ("Delivery Mode", "Instructor-led with lab-based guided practice"),
        ("Assessment", "Open-book written / MCQ assessment on LMS/TMS"),
        ("Target Learners", "Business, operations, data, IT, and cloud learners who need to recognise AI workloads and Azure AI services."),
        ("Prerequisites", "Basic digital literacy and willingness to use the Azure portal, Microsoft Learn, GitHub lab files, and LMS/TMS."),
        ("Trainer Preparation", "Verify slides, Learner Guide, lab repository, LMS/TMS access, practice exam link, Azure accounts/sandbox access, and TRAQOM QR code readiness."),
        ("Course Repository", LAB_REPO),
    ]
    add_doc_table(doc, ["Field", "Detail"], [(k, v) for k, v in info], [1.8, 5.0], break_keywords=())
    note = doc.add_paragraph()
    note.add_run("Note: ").bold = True
    note.add_run("the 'Slide No.' column references the exact slide numbers in the Learner Guide / Course Slides deck for that activity.")

    doc.add_paragraph("Delivery Approach", style="Heading 1")
    doc.add_paragraph(
        "This two-day course is delivered as short concept briefings followed by guided lab practice. "
        "The trainer should repeatedly ask learners to identify the input, required output, Azure AI service, validation evidence, and responsible AI control."
    )
    for item in [
        "Concept briefing: explain the AI workload and service selection principle before each lab.",
        "Guided lab: demonstrate the first few steps, then circulate while learners complete the tasks.",
        "Checkpoint: ask learners to describe the evidence produced and the responsible AI control.",
        "Remediation: revisit misunderstood terms immediately using the relevant slide and lab output.",
        "Closure: confirm attendance, resource cleanup, TRAQOM, and assessment/practice exam actions.",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)

    doc.add_paragraph("Detailed Lesson Schedule", style="Heading 1")
    for day_title, rows in lp_day_rows(slide_index).items():
        doc.add_paragraph(day_title, style="Heading 2")
        add_doc_table(
            doc,
            ["Time", "Lesson Focus", "Trainer Facilitation Notes", "Learner Activity / Evidence", "Slide No."],
            rows,
            [0.75, 1.35, 2.05, 2.05, 0.6],
        )

    doc.add_paragraph("Lab Alignment and Evidence", style="Heading 1")
    doc.add_paragraph(
        "Use this matrix during delivery to confirm that each lab produces evidence linked to the lesson outcome and the slide numbers used in class."
    )
    add_doc_table(
        doc,
        ["Lab", "Lab Title", "Slides", "Trainer Focus", "Learner Evidence / Checkpoint"],
        lp_lab_alignment_rows(labs, slide_index),
        [0.45, 1.65, 0.55, 2.05, 2.10],
        break_keywords=(),
    )

    doc.add_paragraph("Assessment and Funding Requirements", style="Heading 1")
    doc.add_paragraph("The trainer must keep the assessment and funding process visible near the end of Day 2.")
    for item in [
        "Minimum 75% attendance across the course based on SSG digital attendance records.",
        "Completion of assessment digital attendance before the final assessment.",
        "Open-book written / MCQ assessment using approved course materials.",
        "Learners submit answers through LMS/TMS and sign the Assessment Summary Record when instructed.",
        "Completion of feedback and TRAQOM survey as instructed.",
        "Practice exam access: exams.tertiaryinfotech.com for post-course revision.",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item)

    doc.add_paragraph("Trainer Resource Checklist", style="Heading 1")
    add_doc_table(
        doc,
        ["Item", "Trainer Check"],
        [
            ("Slides and Learner Guide", "Confirm latest courseware files are available in the courseware folder and on LMS/TMS."),
            ("Lab repository", f"Confirm learners can access {LAB_REPO} and download or clone lab files."),
            ("Azure access", "Confirm subscription/sandbox availability, portal sign-in, and any instructor-provided credentials."),
            ("Assessment", "Prepare LMS/TMS assessment link, assessment attendance QR, and Assessment Summary Record process."),
            ("Cleanup", "Remind learners to delete or stop personal Azure resources and record any shared cleanup handled by trainer."),
        ],
        [1.55, 5.25],
        break_keywords=(),
    )
    footer(doc)
    doc.settings.element.append(OxmlElement("w:updateFields"))
    doc.save(docx_path)
    return docx_path


def pdf_styles():
    ss = getSampleStyleSheet()
    return {
        "title": ParagraphStyle("Title2", parent=ss["Title"], fontName="Helvetica-Bold", fontSize=22, leading=26, textColor=colors.HexColor("#111827"), alignment=TA_CENTER, spaceAfter=18),
        "h1": ParagraphStyle("H1", parent=ss["Heading1"], fontName="Helvetica-Bold", fontSize=15, leading=18, textColor=colors.HexColor("#1F6FEB"), spaceBefore=10, spaceAfter=6),
        "h2": ParagraphStyle("H2", parent=ss["Heading2"], fontName="Helvetica-Bold", fontSize=12, leading=15, textColor=colors.HexColor("#111827"), spaceBefore=8, spaceAfter=4),
        "body": ParagraphStyle("Body", parent=ss["BodyText"], fontName="Helvetica", fontSize=9.5, leading=12, alignment=TA_LEFT, spaceAfter=5),
        "small": ParagraphStyle("Small", parent=ss["BodyText"], fontName="Helvetica", fontSize=8, leading=10, spaceAfter=4),
        "code": ParagraphStyle("Code", parent=ss["Code"], fontName="Courier", fontSize=7.5, leading=9, backColor=colors.HexColor("#F4F7FB"), borderPadding=4, spaceAfter=6),
    }


def pdf_cover_page(story: list, doc_type: str, st: dict) -> None:
    center_small = ParagraphStyle("CenterSmall", parent=st["body"], alignment=TA_CENTER, fontSize=8, leading=10, textColor=colors.HexColor("#555B66"))
    center_org = ParagraphStyle("CenterOrg", parent=st["body"], alignment=TA_CENTER, fontName="Helvetica-Bold", fontSize=10, leading=12, textColor=colors.HexColor("#111827"))
    center_blue = ParagraphStyle("CenterBlue", parent=st["title"], alignment=TA_CENTER, fontName="Helvetica-Bold", fontSize=24, leading=28, textColor=colors.HexColor("#1F6FEB"))
    center_course = ParagraphStyle("CenterCourse", parent=st["title"], alignment=TA_CENTER, fontName="Helvetica-Bold", fontSize=18, leading=22, textColor=colors.HexColor("#111827"))
    story.append(Spacer(1, 34))
    if TERTIARY_MARK.exists():
        img = RLImage(str(TERTIARY_MARK), width=1.0 * inch, height=1.0 * inch)
        img.hAlign = "CENTER"
        story.append(img)
        story.append(Spacer(1, 4))
    story.append(Paragraph(ORG, center_org))
    story.append(Paragraph("UEN: 201200696W", center_small))
    story.append(Spacer(1, 10))
    story.append(Paragraph(doc_type.upper(), center_blue))
    story.append(Paragraph("For", center_small))
    story.append(Spacer(1, 10))
    if COURSE_BADGE.exists():
        img = RLImage(str(COURSE_BADGE), width=1.05 * inch, height=0.68 * inch)
        img.hAlign = "CENTER"
        story.append(img)
        story.append(Spacer(1, 18))
    story.append(Paragraph(TITLE_WITH_PARENS, center_course))
    story.append(Paragraph(f"TGS Ref No: {CODE}", center_small))
    if "LEARNER" in doc_type.upper():
        story.append(Spacer(1, 8))
        story.append(Paragraph("Conducted by", center_small))
        story.append(Paragraph(ORG, center_org))
        story.append(Paragraph("UEN: 201200696W", center_small))
    story.append(Spacer(1, 8))
    story.append(Paragraph(f"Version {DOC_VERSION}", ParagraphStyle("Version", parent=center_small, textColor=colors.HexColor("#1F6FEB"), fontName="Helvetica-Bold")))
    story.append(PageBreak())


def build_doc_pdf_from_md(md_path: Path, pdf_path: Path, title: str) -> None:
    st = pdf_styles()
    doc = SimpleDocTemplate(str(pdf_path), pagesize=A4, leftMargin=0.65 * inch, rightMargin=0.65 * inch, topMargin=0.65 * inch, bottomMargin=0.65 * inch)
    story = []
    pdf_cover_page(story, "Learner Guide", st)
    story.extend([Paragraph(title, st["title"]), Paragraph(f"{CODE} - Version {DOC_VERSION} - {DATE}", st["body"]), Spacer(1, 8)])
    in_code = False
    code_buf: list[str] = []
    for raw in md_path.read_text(encoding="utf-8").splitlines():
        line = raw.rstrip()
        if line.startswith("```"):
            if in_code:
                story.append(Preformatted("\n".join(code_buf)[:3000], st["code"]))
                code_buf = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue
        if not line:
            continue
        img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
        if img_match:
            img_path = COURSEWARE / img_match.group(2)
            if img_path.exists():
                from PIL import Image
                with Image.open(img_path) as im:
                    iw, ih = im.size
                width = 6.45 * inch
                height = width * ih / iw
                img = RLImage(str(img_path), width=width, height=height)
                img.hAlign = "CENTER"
                story.append(Spacer(1, 5))
                story.append(img)
                story.append(Spacer(1, 8))
            continue
        if line.startswith("# "):
            story.append(Paragraph(strip_md(line[2:]), st["h1"]))
        elif line.startswith("## "):
            story.append(Paragraph(strip_md(line[3:]), st["h1"]))
        elif line.startswith("### "):
            story.append(Paragraph(strip_md(line[4:]), st["h2"]))
        elif line.startswith("- "):
            story.append(ListFlowable([ListItem(Paragraph(strip_md(line[2:]), st["body"]))], bulletType="bullet", leftIndent=16))
        elif re.match(r"^\d+\.\s+", line):
            story.append(Paragraph(strip_md(line), st["body"]))
        elif line.startswith("|"):
            story.append(Preformatted(strip_md(line)[:220], st["small"]))
        else:
            story.append(Paragraph(strip_md(line), st["body"]))
    doc.build(story)


def pdf_cell(value: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(html.escape(str(value)).replace("\n", "<br/>"), style)


def pdf_detail_table(headers: list[str], rows: list[tuple[str, ...]], widths: list[float], st: dict) -> Table:
    data = [[pdf_cell(h, st["small"]) for h in headers]]
    for row in rows:
        data.append([pdf_cell(v, st["small"]) for v in row])
    tbl = Table(data, colWidths=[w * inch for w in widths], repeatRows=1, hAlign="CENTER")
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F6FEB")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#DADCE0")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    for idx, row in enumerate(rows, 1):
        if "break" in " ".join(row).lower():
            tbl.setStyle(TableStyle([("BACKGROUND", (0, idx), (-1, idx), colors.HexColor("#FFF4E5"))]))
    return tbl


def build_lp_pdf(labs: list[Lab], slide_index: list[dict[str, str | int]], pdf_path: Path) -> None:
    st = pdf_styles()
    doc = SimpleDocTemplate(str(pdf_path), pagesize=A4, leftMargin=0.55 * inch, rightMargin=0.55 * inch, topMargin=0.55 * inch, bottomMargin=0.55 * inch)
    story = []
    pdf_cover_page(story, "Lesson Plan", st)
    story.extend([
        Paragraph(f"Lesson Plan - {TITLE_WITH_PARENS}", st["title"]),
        Paragraph(f"{CODE} - Version {DOC_VERSION} - {DATE}", st["body"]),
        Paragraph("Course Information", st["h1"]),
    ])
    story.append(pdf_detail_table(
        ["Field", "Detail"],
        [
            ("Duration", "2 Days / 16 training hours, 9:00 AM - 6:00 PM"),
            ("Delivery Mode", "Instructor-led with lab-based guided practice"),
            ("Assessment", "Open-book written / MCQ assessment on LMS/TMS"),
            ("Trainer Preparation", "Verify slides, Learner Guide, lab repository, LMS/TMS, Azure access, practice exam link, and TRAQOM QR codes."),
            ("Course Repository", LAB_REPO),
        ],
        [1.45, 5.0],
        st,
    ))
    story.append(Paragraph("Delivery Approach", st["h1"]))
    for item in [
        "Use short concept briefings followed by guided lab practice.",
        "For each lab, ask learners to identify input, output, Azure service, validation evidence, and responsible AI control.",
        "Use checkpoints for remediation before moving to the next topic.",
        "Keep attendance, assessment, TRAQOM, cleanup, and practice exam actions visible near the end of Day 2.",
    ]:
        story.append(ListFlowable([ListItem(Paragraph(item, st["body"]))], bulletType="bullet", leftIndent=16))

    story.append(PageBreak())
    story.append(Paragraph("Detailed Lesson Schedule", st["h1"]))
    for day_title, rows in lp_day_rows(slide_index).items():
        story.append(Paragraph(day_title, st["h2"]))
        story.append(pdf_detail_table(
            ["Time", "Lesson Focus", "Trainer Notes", "Learner Evidence", "Slides"],
            rows,
            [0.62, 1.1, 1.85, 1.85, 0.65],
            st,
        ))
        story.append(Spacer(1, 8))

    story.append(PageBreak())
    story.append(Paragraph("Lab Alignment and Evidence", st["h1"]))
    story.append(Paragraph("Use this matrix to confirm every lab produces evidence linked to the lesson outcome and slide numbers.", st["body"]))
    story.append(pdf_detail_table(
        ["Lab", "Lab Title", "Slides", "Trainer Focus", "Learner Evidence / Checkpoint"],
        lp_lab_alignment_rows(labs, slide_index),
        [0.38, 1.32, 0.50, 1.85, 2.0],
        st,
    ))
    story.append(PageBreak())
    story.append(Paragraph("Assessment and Funding Requirements", st["h1"]))
    for item in [
        "Minimum 75% attendance across the course based on SSG digital attendance records.",
        "Complete assessment digital attendance before final assessment.",
        "Use approved course materials only for the open-book written / MCQ assessment.",
        "Submit answers through LMS/TMS and sign the Assessment Summary Record when instructed.",
        "Complete feedback and TRAQOM survey; use exams.tertiaryinfotech.com for revision.",
    ]:
        story.append(ListFlowable([ListItem(Paragraph(item, st["body"]))], bulletType="bullet", leftIndent=16))
    story.append(Paragraph("Trainer Resource Checklist", st["h1"]))
    story.append(pdf_detail_table(
        ["Item", "Trainer Check"],
        [
            ("Slides and Learner Guide", "Confirm latest courseware files are in the courseware folder and on LMS/TMS."),
            ("Lab repository", f"Confirm learners can access {LAB_REPO}."),
            ("Azure access", "Confirm subscription/sandbox availability and portal sign-in."),
            ("Assessment", "Prepare LMS/TMS assessment link, attendance QR, and Assessment Summary Record process."),
            ("Cleanup", "Remind learners to delete or stop personal Azure resources."),
        ],
        [1.4, 5.05],
        st,
    ))
    doc.build(story)


def build_slide_pdf(slide_index: list[dict[str, str | int]], labs: list[Lab], pdf_path: Path) -> None:
    st = pdf_styles()
    doc = SimpleDocTemplate(str(pdf_path), pagesize=landscape(A4), leftMargin=0.55 * inch, rightMargin=0.55 * inch, topMargin=0.45 * inch, bottomMargin=0.45 * inch)
    story = []
    for row in slide_index:
        title = str(row["title"])
        section = str(row["section"])
        story.append(Paragraph(title, st["title"]))
        story.append(Paragraph(f"Slide {row['slide']} - {section}", st["body"]))
        if section.startswith("Lab "):
            n = int(section.split()[1])
            lab = next((x for x in labs if x.num == n), None)
            if lab:
                story.append(Paragraph(lab.scenario, st["body"]))
                for item in lab.objectives[:5]:
                    story.append(ListFlowable([ListItem(Paragraph(item, st["body"]))], bulletType="bullet", leftIndent=18))
                story.append(Paragraph(f"Validation: {lab.validation}", st["body"]))
        else:
            story.append(Paragraph(f"{TITLE_WITH_PARENS} course slide generated from the lab-aligned courseware source.", st["body"]))
        story.append(PageBreak())
    doc.build(story)


def main() -> None:
    ensure_dirs()
    extract_reference_media()
    labs = parse_labs()
    if len(labs) != 10:
        raise RuntimeError(f"Expected 10 labs, found {len(labs)}")
    pptx_path, slide_index = build_slides(labs)
    md_path, lg_docx = build_lg(labs)
    lp_docx = build_lp(labs, slide_index)
    build_slide_pdf(slide_index, labs, COURSEWARE / f"{TITLE}-{VERSION}.pdf")
    build_doc_pdf_from_md(md_path, COURSEWARE / f"LG-{TITLE}.pdf", f"Learner Guide - {TITLE_WITH_PARENS}")
    build_lp_pdf(labs, slide_index, COURSEWARE / f"LP-{TITLE}.pdf")
    (COURSEWARE / "slide-index.txt").write_text(
        "\n".join(f"{r['slide']:>3}. [{r['section']}] {r['title']}" for r in slide_index),
        encoding="utf-8",
    )
    print(f"Generated {pptx_path}")
    print(f"Generated {COURSEWARE / (TITLE + '-' + VERSION + '.pdf')}")
    print(f"Generated {lg_docx}")
    print(f"Generated {COURSEWARE / ('LG-' + TITLE + '.pdf')}")
    print(f"Generated {md_path}")
    print(f"Generated {lp_docx}")
    print(f"Generated {COURSEWARE / ('LP-' + TITLE + '.pdf')}")
    print(f"Slides: {len(slide_index)}")


if __name__ == "__main__":
    main()
